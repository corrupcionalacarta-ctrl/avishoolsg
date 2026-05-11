"""
drive_analyzer.py - Análisis de archivos Drive con Gemini
==========================================================
Procesa los archivos descargados de Drive "Compartido conmigo" con Gemini:
  - Extrae texto y contenido (multimodal para PDFs/imágenes)
  - Clasifica el contenido: prueba, guía, pauta, material
  - Extrae temas, preguntas, ejercicios y conceptos clave
  - Genera un resumen para el tutor IA
  - Persiste en Supabase tabla `classroom_archivos`

Uso:
    python drive_analyzer.py              # analiza todos los alumnos
    python drive_analyzer.py clemente     # solo Clemente
    python drive_analyzer.py --force      # re-analizar aunque ya estén analizados
"""

import base64
import json
import os
import re
import sys
from pathlib import Path

from dotenv import load_dotenv

load_dotenv()

OUTPUT_DIR   = Path(os.getenv("OUTPUT_DIR", "."))
DOWNLOAD_DIR = OUTPUT_DIR / "drive_downloads"
GEMINI_KEY   = (os.getenv("GEMINI_API_KEY") or "").strip()
GEMINI_MODEL = (os.getenv("GEMINI_MODEL") or "gemini-2.5-flash").strip()

ALUMNOS_CONFIG = []
for i in (1, 2):
    nombre = (os.getenv(f"ALUMNO_{i}_NOMBRE") or "").strip()
    email  = (os.getenv(f"ALUMNO_{i}_CLASSROOM") or "").strip()
    if nombre and email:
        ALUMNOS_CONFIG.append({
            "nombre": nombre,
            "slug":   nombre.split()[0].lower(),
        })

# Tipos de archivo que Gemini puede procesar directamente
GEMINI_MIME = {
    ".pdf":  "application/pdf",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png":  "image/png",
}

MAX_FILE_SIZE_MB = 15  # Gemini inline limit


# ─────────────────────────────────────────────────────────────────────────────
# Gemini helpers
# ─────────────────────────────────────────────────────────────────────────────

def get_gemini_client():
    if not GEMINI_KEY:
        raise ValueError("GEMINI_API_KEY no configurada en .env")
    from google import genai
    return genai.Client(api_key=GEMINI_KEY)


def _safe_print(msg: str):
    sys.stdout.buffer.write((msg + "\n").encode("utf-8", "replace"))
    sys.stdout.buffer.flush()


def _extract_pptx_text(file_path: Path) -> str:
    """Extrae texto de un PPTX slide a slide."""
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides_text.append(f"[Diapositiva {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides_text)
    except Exception as e:
        _safe_print(f"  [WARN] No se pudo extraer texto de PPTX: {e}")
        return ""


def _build_prompt(file_name: str, alumno_nombre: str, curso_hint: str = "", extra_context: str = "") -> str:
    return f"""Analiza este documento educativo de un alumno de 6to básico en Chile ({alumno_nombre}).

El archivo se llama: "{file_name}"
{f'Asignatura probable: {curso_hint}' if curso_hint else ''}
{extra_context}

Responde SOLO con un JSON válido con este esquema exacto (sin texto extra, sin markdown):
{{
  "tipo_contenido": "<prueba|guia|pauta|material|ejercicio|temario|otro>",
  "asignatura": "<Matemática|Lenguaje|Ciencias|Historia|Inglés|Ed.Física|Arte|Música|Tecnología|Filosofía|Teología|Orientación|otro>",
  "titulo_inferido": "<título descriptivo del documento>",
  "unidad_tematica": "<unidad o tema principal>",
  "temas": ["<tema 1>", "<tema 2>", "..."],
  "conceptos_clave": ["<concepto 1>", "<concepto 2>", "..."],
  "preguntas_o_ejercicios": [
    {{"numero": 1, "enunciado": "<enunciado breve>", "tipo": "<seleccion_multiple|desarrollo|completar|verdadero_falso|calculo|otro>"}},
    ...
  ],
  "nivel_dificultad": "<facil|medio|dificil>",
  "resumen": "<2-3 oraciones describiendo el contenido para el tutor>",
  "fecha_probable": "<YYYY-MM si se infiere del nombre, null si no>",
  "tiene_respuestas": <true si es pauta/solucionario, false si no>
}}

Importante:
- Si es una prueba, extrae TODAS las preguntas con su enunciado
- Si es una guía, extrae todos los ejercicios
- Si es material/presentación del profe (PPT), resume los temas cubiertos en cada diapositiva
- Si es una pauta, extrae las respuestas correctas en el campo preguntas
- El campo "resumen" debe ser útil para que el tutor sepa de qué trata el documento
"""


def _parse_gemini_json(text: str) -> dict | None:
    text = text.strip()
    json_match = re.search(r'\{.*\}', text, re.DOTALL)
    if json_match:
        text = json_match.group(0)
    try:
        return json.loads(text)
    except Exception:
        return None


def analyze_file_with_gemini(client, file_path: Path, alumno_nombre: str, curso_hint: str = "") -> dict | None:
    """
    Analiza un archivo con Gemini.
    - PDFs y DOCX: envío inline como bytes (multimodal)
    - PPTX: extrae texto con python-pptx, envía como texto
    Retorna dict con: tipo_contenido, asignatura, temas, resumen, preguntas, nivel_dificultad
    """
    from google import genai
    from google.genai import types

    suffix = file_path.suffix.lower()
    size_mb = file_path.stat().st_size / (1024 * 1024)

    # ── PPTX: extraer texto y enviar como texto plano ──────────────────────────
    if suffix in (".pptx", ".ppt"):
        pptx_text = _extract_pptx_text(file_path)
        if not pptx_text.strip():
            _safe_print(f"  [SKIP] {file_path.name} — PPTX sin texto extraíble")
            return None

        prompt = _build_prompt(file_path.name, alumno_nombre, curso_hint,
                               extra_context=f"\nContenido de la presentación:\n{pptx_text[:12000]}")
        try:
            response = client.models.generate_content(
                model=GEMINI_MODEL,
                contents=[types.Part.from_text(text=prompt)],
            )
            return _parse_gemini_json(response.text)
        except Exception as e:
            _safe_print(f"  [ERROR] Gemini PPTX falló en {file_path.name}: {e}")
            return None

    # ── Archivos binarios: enviar inline como bytes ────────────────────────────
    mime_type = GEMINI_MIME.get(suffix)
    if not mime_type:
        _safe_print(f"  [SKIP] {file_path.name} — tipo no soportado ({suffix})")
        return None

    if size_mb > MAX_FILE_SIZE_MB:
        _safe_print(f"  [SKIP] {file_path.name} — muy grande ({size_mb:.1f}MB > {MAX_FILE_SIZE_MB}MB)")
        return None

    file_bytes = file_path.read_bytes()
    prompt = _build_prompt(file_path.name, alumno_nombre, curso_hint)

    try:
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[
                types.Part.from_bytes(data=file_bytes, mime_type=mime_type),
                types.Part.from_text(text=prompt),
            ],
        )
        return _parse_gemini_json(response.text)
    except Exception as e:
        _safe_print(f"  [ERROR] Gemini falló en {file_path.name}: {e}")
        return None


# ─────────────────────────────────────────────────────────────────────────────
# Supabase — tabla classroom_archivos
# ─────────────────────────────────────────────────────────────────────────────

def ensure_table(sb):
    """Crea la tabla classroom_archivos si no existe."""
    sql = """
    CREATE TABLE IF NOT EXISTS classroom_archivos (
        id              uuid PRIMARY KEY DEFAULT gen_random_uuid(),
        alumno          text NOT NULL,
        archivo_nombre  text NOT NULL,
        archivo_path    text,
        asignatura      text,
        tipo_contenido  text,
        titulo_inferido text,
        unidad_tematica text,
        temas           jsonb DEFAULT '[]',
        conceptos_clave jsonb DEFAULT '[]',
        preguntas       jsonb DEFAULT '[]',
        nivel_dificultad text,
        resumen         text,
        fecha_probable  text,
        tiene_respuestas boolean DEFAULT false,
        analizado_en    timestamptz DEFAULT now(),
        UNIQUE (alumno, archivo_nombre)
    );
    CREATE INDEX IF NOT EXISTS idx_clasarch_alumno     ON classroom_archivos (alumno);
    CREATE INDEX IF NOT EXISTS idx_clasarch_asignatura ON classroom_archivos (asignatura);
    CREATE INDEX IF NOT EXISTS idx_clasarch_tipo       ON classroom_archivos (tipo_contenido);
    """
    try:
        sb.rpc("exec_ddl", {"sql": sql}).execute()
    except Exception:
        pass  # La tabla puede ya existir o no tener la RPC disponible


def push_analysis(alumno_nombre: str, archivo_nombre: str, archivo_path: str, analysis: dict) -> bool:
    """Upsert del análisis en classroom_archivos."""
    url = (os.getenv("SUPABASE_URL") or "").strip()
    key = (os.getenv("SUPABASE_SERVICE_KEY") or "").strip()
    if not url or not key:
        return False

    try:
        from supabase import create_client
        sb = create_client(url, key)

        row = {
            "alumno":           alumno_nombre,
            "archivo_nombre":   archivo_nombre[:300],
            "archivo_path":     archivo_path[:500],
            "asignatura":       analysis.get("asignatura", "")[:100],
            "tipo_contenido":   analysis.get("tipo_contenido", "material")[:50],
            "titulo_inferido":  (analysis.get("titulo_inferido") or "")[:300],
            "unidad_tematica":  (analysis.get("unidad_tematica") or "")[:300],
            "temas":            json.dumps(analysis.get("temas", []), ensure_ascii=False),
            "conceptos_clave":  json.dumps(analysis.get("conceptos_clave", []), ensure_ascii=False),
            "preguntas":        json.dumps(analysis.get("preguntas_o_ejercicios", []), ensure_ascii=False),
            "nivel_dificultad": (analysis.get("nivel_dificultad") or "medio")[:20],
            "resumen":          (analysis.get("resumen") or "")[:1000],
            "fecha_probable":   analysis.get("fecha_probable"),
            "tiene_respuestas": bool(analysis.get("tiene_respuestas", False)),
        }

        # Upsert por (alumno, archivo_nombre)
        sb.table("classroom_archivos").upsert(row, on_conflict="alumno,archivo_nombre").execute()
        return True

    except Exception as e:
        err = str(e)
        if "PGRST205" not in err:  # Silenciar "tabla no existe" — se guarda en local
            _safe_print(f"  [WARN] Supabase upsert: {err[:120]}")
        return False


def get_already_analyzed(alumno_nombre: str) -> set[str]:
    """Retorna set de nombres de archivos ya analizados en Supabase."""
    url = (os.getenv("SUPABASE_URL") or "").strip()
    key = (os.getenv("SUPABASE_SERVICE_KEY") or "").strip()
    if not url or not key:
        return set()

    try:
        from supabase import create_client
        sb = create_client(url, key)
        resp = sb.table("classroom_archivos").select("archivo_nombre").eq("alumno", alumno_nombre).execute()
        return {row["archivo_nombre"] for row in (resp.data or [])}
    except Exception:
        return set()


# ─────────────────────────────────────────────────────────────────────────────
# Pipeline principal
# ─────────────────────────────────────────────────────────────────────────────

# Inferir asignatura desde el nombre del archivo
ASIGNATURA_HINTS = {
    "matemát": "Matemática", "matematica": "Matemática", "mcm": "Matemática",
    "mínimo": "Matemática", "minimo": "Matemática", "múltiplo": "Matemática",
    "multiplo": "Matemática", "fraccion": "Matemática", "fracción": "Matemática",
    "número": "Matemática", "numero": "Matemática", "primo": "Matemática",
    "divisor": "Matemática", "factor": "Matemática", "gráfico": "Matemática",
    "grafico": "Matemática", "probabilidad": "Matemática", "tallo": "Matemática",
    "barra": "Matemática",
    "fotosíntesis": "Ciencias", "fotosintesis": "Ciencias", "planta": "Ciencias",
    "hidrosfera": "Ciencias", "hidrósfera": "Ciencias", "microscopio": "Ciencias",
    "microscopía": "Ciencias", "lab": "Ciencias", "laboratorio": "Ciencias",
    "agua": "Ciencias", "estoma": "Ciencias", "cloroplasto": "Ciencias",
    "independencia": "Historia", "historia": "Historia",
    "listening": "Inglés", "reading": "Inglés", "writing": "Inglés",
    "speaking": "Inglés", "english": "Inglés", "bimester": "Inglés",
    "character": "Inglés", "answer": "Inglés",
}


def infer_asignatura(filename: str) -> str:
    name_lower = filename.lower()
    for keyword, asig in ASIGNATURA_HINTS.items():
        if keyword in name_lower:
            return asig
    return ""


def analyze_alumno(alumno: dict, force: bool = False):
    """Analiza todos los archivos descargados de un alumno."""
    slug = alumno["slug"]
    nombre = alumno["nombre"]
    download_dir = DOWNLOAD_DIR / slug

    if not download_dir.exists():
        _safe_print(f"[WARN] No hay archivos descargados para {nombre} en {download_dir}")
        return

    _safe_print(f"\n{'='*60}")
    _safe_print(f"[Analyzer] {nombre}")
    _safe_print(f"{'='*60}")

    if not GEMINI_KEY:
        _safe_print("[ERROR] GEMINI_API_KEY no configurada")
        return

    try:
        client = get_gemini_client()
    except Exception as e:
        _safe_print(f"[ERROR] Gemini client: {e}")
        return

    # Archivos ya analizados (para skip)
    already_done = set() if force else get_already_analyzed(nombre)
    _safe_print(f"[INFO] {len(already_done)} archivos ya analizados previamente")

    # Listar archivos descargados
    files = sorted(download_dir.iterdir(), key=lambda p: p.name.lower())
    _safe_print(f"[INFO] {len(files)} archivos en {download_dir}")

    ok = 0
    skip = 0
    errors = 0

    for file_path in files:
        if not file_path.is_file():
            continue

        if file_path.name in already_done:
            skip += 1
            continue

        suffix = file_path.suffix.lower()
        if suffix not in GEMINI_MIME:
            _safe_print(f"  [SKIP] {file_path.name} — formato no soportado")
            skip += 1
            continue

        _safe_print(f"  [→] {file_path.name}")

        asig_hint = infer_asignatura(file_path.name)
        analysis = analyze_file_with_gemini(client, file_path, nombre, asig_hint)

        if not analysis:
            errors += 1
            continue

        # Mostrar resumen del análisis
        tipo = analysis.get("tipo_contenido", "?")
        asig = analysis.get("asignatura", "?")
        n_preguntas = len(analysis.get("preguntas_o_ejercicios", []))
        resumen_short = (analysis.get("resumen") or "")[:80]
        _safe_print(f"     [{tipo.upper():<8}] [{asig:<12}] {n_preguntas} preguntas | {resumen_short}")

        # Persistir en Supabase
        # Guardar siempre en local (rápido, no depende de Supabase)
        local_dir = OUTPUT_DIR / f"analysis_{slug}"
        local_dir.mkdir(exist_ok=True)
        local_path = local_dir / f"{file_path.stem[:80]}.json"
        local_path.write_text(
            json.dumps({"archivo": file_path.name, "analisis": analysis}, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

        # También intentar Supabase (puede fallar si la tabla no existe aún)
        push_analysis(nombre, file_path.name, str(file_path), analysis)
        ok += 1

    _safe_print(f"\n[RESUMEN] {nombre}: {ok} analizados, {skip} saltados, {errors} errores")

    # Guardar resumen consolidado
    save_consolidated_analysis(nombre, slug)


def save_consolidated_analysis(alumno_nombre: str, slug: str):
    """Guarda un resumen consolidado de todos los análisis en JSON local."""
    url = (os.getenv("SUPABASE_URL") or "").strip()
    key = (os.getenv("SUPABASE_SERVICE_KEY") or "").strip()
    if not url or not key:
        return

    try:
        from supabase import create_client
        sb = create_client(url, key)
        resp = sb.table("classroom_archivos").select("*").eq("alumno", alumno_nombre).execute()
        data = resp.data or []

        out_path = OUTPUT_DIR / f"analysis_summary_{slug}.json"
        out_path.write_text(
            json.dumps(data, ensure_ascii=False, indent=2, default=str),
            encoding="utf-8"
        )
        _safe_print(f"[DEBUG] Resumen consolidado en {out_path} ({len(data)} registros)")

    except Exception as e:
        if "PGRST205" not in str(e):
            _safe_print(f"[WARN] No se pudo guardar resumen: {e}")


def main():
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("alumno", nargs="?", help="Slug del alumno (ej: clemente)")
    parser.add_argument("--force", action="store_true", help="Re-analizar aunque ya estén en Supabase")
    args = parser.parse_args()

    alumnos = ALUMNOS_CONFIG
    if args.alumno:
        alumnos = [a for a in alumnos if a["slug"] == args.alumno.lower()]

    if not alumnos:
        _safe_print("[ERROR] No hay alumnos configurados o el slug no existe")
        sys.exit(1)

    for alumno in alumnos:
        analyze_alumno(alumno, force=args.force)


if __name__ == "__main__":
    main()
